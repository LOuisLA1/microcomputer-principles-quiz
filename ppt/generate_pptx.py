#!/usr/bin/env python3
"""Generate an editable .pptx from the Claude Code vs Codex HTML PPT content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (IKB Swiss) ──
IKB       = RGBColor(0x00, 0x2F, 0xA7)  # #002FA7
IKB_BRIGHT = RGBColor(0x5B, 0x7B, 0xFF)
INK       = RGBColor(0x0A, 0x0A, 0x0A)  # near black
PAPER     = RGBColor(0xFA, 0xFA, 0xF8)  # warm white
GREY_1    = RGBColor(0xF0, 0xF0, 0xEE)
GREY_2    = RGBColor(0xD4, 0xD4, 0xD2)
GREY_3    = RGBColor(0x73, 0x73, 0x73)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="", font_size=Pt(12),
                font_name="Microsoft YaHei", color=INK, bold=False, alignment=PP_ALIGN.LEFT,
                line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = line_spacing
    return tf


def add_rich_textbox(slide, left, top, width, height):
    """Add a textbox and return the text_frame for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, font_size=Pt(12), font_name="Microsoft YaHei", color=INK,
             bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.2, space_after=Pt(0)):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = space_after
    p.line_spacing = line_spacing
    return p


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_line(slide, left, top, width, height, color=GREY_2, width_pt=Pt(1)):
    """Add a horizontal line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def chrome_text(slide, left_line, right_line, page_n):
    """Add chrome header: left text | page number."""
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(8), Inches(0.35),
                left_line, Pt(9), "Consolas", GREY_3, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(8), Inches(0.35), Inches(5), Inches(0.35),
                f"{page_n} / 10", Pt(9), "Consolas", GREY_3, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 · S01 Cover · IKB Background
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, IKB)
chrome_text(s, "Internal Tech Talk · 内部技术分享", "01 / 10", 1)

add_textbox(s, Inches(0.6), Inches(1.4), Inches(10), Inches(0.4),
            "AI CODING TOOLS · COMPARISON", Pt(11), "Consolas",
            RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.LEFT)

add_textbox(s, Inches(0.6), Inches(2.2), Inches(12), Inches(2.8),
            "Claude Code\nvs Codex", Pt(72), "Microsoft YaHei",
            WHITE, alignment=PP_ALIGN.LEFT, line_spacing=0.95)

add_line(s, Inches(0.6), Inches(5.6), Inches(8), Pt(1),
         RGBColor(0xFF, 0xFF, 0xFF))

add_textbox(s, Inches(0.6), Inches(5.8), Inches(8), Inches(0.8),
            "两个 AI 编程工具的架构、哲学与工作流深度对比 · 帮助团队做出正确的选型决策",
            Pt(16), "Microsoft YaHei", RGBColor(0xDD, 0xDD, 0xDD),
            alignment=PP_ALIGN.LEFT, line_spacing=1.4)

add_textbox(s, Inches(0.6), Inches(6.7), Inches(6), Inches(0.3),
            "内部技术分享 · 2026.05", Pt(10), "Consolas",
            RGBColor(0xAA, 0xAA, 0xAA))

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 · S08 Duo Compare
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, PAPER)
chrome_text(s, "Overview · 基本信息", "02 / 10", 2)

# Left column - Claude Code
add_textbox(s, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.3),
            "CLAUDE CODE", Pt(10), "Consolas", IKB, bold=True)
add_textbox(s, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.6),
            "Anthropic", Pt(36), "Microsoft YaHei", INK, line_spacing=1.0)

items_cc = [
    ("模型", "Claude Opus 4.7 / Sonnet 4.6 / Haiku 4.5"),
    ("运行环境", "本地 CLI · VS Code / JetBrains IDE"),
    ("发布时间", "2024 年 · 持续迭代至今"),
    ("核心哲学", "安全优先 · 深度推理 · 人机协作"),
]
y = 2.4
for label, desc in items_cc:
    add_textbox(s, Inches(0.8), Inches(y), Inches(5.2), Inches(0.25),
                label, Pt(10), "Consolas", IKB, bold=True)
    add_textbox(s, Inches(0.8), Inches(y + 0.3), Inches(5.2), Inches(0.3),
                desc, Pt(12), "Microsoft YaHei", GREY_3)
    y += 0.8

# Vertical divider
add_line(s, Inches(6.4), Inches(1.2), Pt(1), Inches(5.5), GREY_2)

# Right column - Codex
add_textbox(s, Inches(6.8), Inches(1.1), Inches(5.5), Inches(0.3),
            "CODEX", Pt(10), "Consolas", IKB, bold=True)
add_textbox(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.6),
            "OpenAI", Pt(36), "Microsoft YaHei", INK, line_spacing=1.0)

items_cx = [
    ("模型", "GPT-4o / GPT-4.1 · o3 / o4-mini"),
    ("运行环境", "CLI · Cloud Sandbox 云端执行"),
    ("发布时间", "2024 年 · 持续迭代至今"),
    ("核心哲学", "广度覆盖 · 多模态 · 自主执行"),
]
y = 2.4
for label, desc in items_cx:
    add_textbox(s, Inches(7.0), Inches(y), Inches(5.2), Inches(0.25),
                label, Pt(10), "Consolas", IKB, bold=True)
    add_textbox(s, Inches(7.0), Inches(y + 0.3), Inches(5.2), Inches(0.3),
                desc, Pt(12), "Microsoft YaHei", GREY_3)
    y += 0.8

add_textbox(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.3),
            "同一赛道 · 两种路线 · 理解差异才能做出正确选择",
            Pt(11), "Consolas", GREY_3, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 · S03 Statement (Dark)
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, INK)
chrome_text(s, "Key Insight · 核心洞察", "03 / 10", 3)

tf = add_rich_textbox(s, Inches(1.5), Inches(2.0), Inches(10), Inches(4.0))
add_para(tf, "同样让 AI", Pt(48), "Microsoft YaHei", WHITE, line_spacing=1.1)
add_para(tf, "写代码,", Pt(48), "Microsoft YaHei", IKB_BRIGHT, line_spacing=1.1)
add_para(tf, "底层的选择,", Pt(48), "Microsoft YaHei", IKB_BRIGHT, line_spacing=1.1)
add_para(tf, "决定了你的", Pt(48), "Microsoft YaHei", WHITE, line_spacing=1.1)
add_para(tf, "工作流.", Pt(48), "Microsoft YaHei", IKB_BRIGHT, line_spacing=1.1)

add_textbox(s, Inches(1.5), Inches(6.2), Inches(10), Inches(0.3),
            "— The tools shape the workflow, not the other way around.",
            Pt(11), "Consolas", RGBColor(0x88, 0x88, 0x88))

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 · S19 Four Cards
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, PAPER)
chrome_text(s, "Dimensions · 对比维度", "04 / 10", 4)

add_line(s, Inches(0.6), Inches(1.1), Inches(1.2), Pt(2), IKB)
add_textbox(s, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
            "四大维度\n深度对比", Pt(36), "Microsoft YaHei", INK,
            line_spacing=0.95)

cards = [
    ("01 / MODEL", "模型策略",
     "Claude: 安全对齐 + 深度推理\nGPT: 广度覆盖 + 多模态能力\n不同的训练哲学导致不同的输出风格"),
    ("02 / EXTEND", "扩展体系",
     "Claude: Skills + Hooks + MCP\nCodex: Sandbox + Plugins\n声明式配置 vs 编程式扩展"),
    ("03 / WORKFLOW", "工作流哲学",
     "Claude: Plan → Task → Memory\nCodex: Task → Agent Loop\n显式规划 vs 自主探索"),
    ("04 / RUNTIME", "运行环境",
     "Claude: 本地终端直接执行\nCodex: 云端沙箱隔离运行\n本地控制 vs 云端安全"),
]

positions = [
    (Inches(0.6), Inches(2.9)),
    (Inches(6.5), Inches(2.9)),
    (Inches(0.6), Inches(5.1)),
    (Inches(6.5), Inches(5.1)),
]

for (tag, title, body), (x, y) in zip(cards, positions):
    add_rect(s, x, y, Inches(5.5), Inches(1.8), fill_color=GREY_1)
    add_textbox(s, x + Inches(0.3), y + Inches(0.2), Inches(4.9), Inches(0.25),
                tag, Pt(9), "Consolas", GREY_3)
    add_textbox(s, x + Inches(0.3), y + Inches(0.5), Inches(4.9), Inches(0.4),
                title, Pt(20), "Microsoft YaHei", INK)
    add_textbox(s, x + Inches(0.3), y + Inches(1.0), Inches(4.9), Inches(0.7),
                body, Pt(11), "Microsoft YaHei", GREY_3, line_spacing=1.4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 · S11 Horizontal Timeline
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, GREY_1)
chrome_text(s, "Evolution · 演进历程", "05 / 10", 5)

add_textbox(s, Inches(0.6), Inches(1.0), Inches(6), Inches(0.25),
            "TIMELINE", Pt(10), "Consolas", IKB, bold=True)
add_textbox(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.8),
            "两个工具的演进之路", Pt(36), "Microsoft YaHei", INK,
            line_spacing=0.95)

# Timeline axis
add_line(s, Inches(0.8), Inches(4.2), Inches(11.5), Pt(1), GREY_2)

timeline = [
    ("2024 Q1", "Claude Code", "初始发布"),
    ("2024 Q4", "Skills & Hooks", "可编程扩展"),
    ("2025 Q2", "Plan + MCP", "显式规划"),
    ("2025 Q4", "Memory + Worktree", "持久化隔离"),
    ("2026", "Agent Ecosystem", "技能市场"),
]

for i, (yr, name, desc) in enumerate(timeline):
    x = Inches(0.8 + i * 2.4)
    # Dot
    add_rect(s, x + Inches(0.9), Inches(4.1), Inches(0.15), Inches(0.15),
             fill_color=IKB)
    # Label above
    if i % 2 == 0:
        add_textbox(s, x, Inches(3.0), Inches(2.0), Inches(0.25),
                    yr, Pt(9), "Consolas", GREY_3, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(3.25), Inches(2.0), Inches(0.3),
                    name, Pt(12), "Microsoft YaHei", INK, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(3.5), Inches(2.0), Inches(0.25),
                    desc, Pt(10), "Microsoft YaHei", GREY_3, alignment=PP_ALIGN.CENTER)
    else:
        add_textbox(s, x, Inches(4.5), Inches(2.0), Inches(0.25),
                    yr, Pt(9), "Consolas", GREY_3, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(4.75), Inches(2.0), Inches(0.3),
                    name, Pt(12), "Microsoft YaHei", INK, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(5.0), Inches(2.0), Inches(0.25),
                    desc, Pt(10), "Microsoft YaHei", GREY_3, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 · S04 Six Cells (Dark)
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, INK)
chrome_text(s, "Key Differences · 关键差异", "06 / 10", 6)

add_textbox(s, Inches(0.6), Inches(1.0), Inches(6), Inches(0.25),
            "SIX DIMENSIONS", Pt(10), "Consolas", IKB_BRIGHT, bold=True)

cells = [
    ("01", "模型底座", "Claude Opus/Sonnet/Haiku vs GPT-4o/o3/o4-mini"),
    ("02", "扩展机制", "Skills 声明式 + Hooks 事件驱动 vs Plugins 编程式"),
    ("03", "规划能力", "Plan Mode 显式架构设计 vs Agent Loop 自主分解"),
    ("04", "上下文持久化", "Memory 跨会话记忆系统 vs Sandbox 会话级状态"),
    ("05", "安全模型", "权限分级 + 用户确认 vs 沙箱隔离 + 自动审批"),
    ("06", "工作隔离", "Git Worktree 原生支持 vs Sandbox Branch 虚拟分支"),
]

positions = [
    (Inches(0.6), Inches(1.6)), (Inches(4.5), Inches(1.6)), (Inches(8.4), Inches(1.6)),
    (Inches(0.6), Inches(3.8)), (Inches(4.5), Inches(3.8)), (Inches(8.4), Inches(3.8)),
]

for (num, title, desc), (x, y) in zip(cells, positions):
    is_accent = (num == "06")
    bg = IKB if is_accent else RGBColor(0x22, 0x22, 0x22)
    text_c = WHITE if is_accent else RGBColor(0xCC, 0xCC, 0xCC)
    add_rect(s, x, y, Inches(3.5), Inches(1.9), fill_color=bg)
    add_textbox(s, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.3),
                num, Pt(11), "Consolas", RGBColor(0x88, 0x88, 0x88))
    add_textbox(s, x + Inches(0.25), y + Inches(0.6), Inches(3.0), Inches(0.35),
                title, Pt(16), "Microsoft YaHei", text_c)
    add_textbox(s, x + Inches(0.25), y + Inches(1.1), Inches(3.0), Inches(0.6),
                desc, Pt(10), "Microsoft YaHei", RGBColor(0x99, 0x99, 0x99),
                line_spacing=1.4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 · S05 Three Layers
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, PAPER)
chrome_text(s, "Architecture · 架构对比", "07 / 10", 7)

add_textbox(s, Inches(0.6), Inches(1.0), Inches(6), Inches(0.25),
            "THREE LAYERS", Pt(10), "Consolas", IKB, bold=True)
add_textbox(s, Inches(0.6), Inches(1.3), Inches(12), Inches(1.0),
            "两者都分三层\n但每一层都不同", Pt(36), "Microsoft YaHei", INK,
            line_spacing=0.95)

layers = [
    ("01", "交互层 · Interface",
     "Claude Code: CLI + IDE 集成 · 文件级操作 · 实时终端反馈\nCodex: CLI + Web UI · 沙箱级操作 · 异步状态推送"),
    ("02", "编排层 · Orchestration",
     "Claude Code: Plan Mode 显式规划 · Task 跟踪 · Agent 子代理分发\nCodex: Agent Loop 自动规划 · Sandbox 任务队列 · 异步多步执行"),
    ("03", "扩展层 · Extension",
     "Claude Code: Skills + Hooks + MCP · 文件即技能 · 事件驱动\nCodex: Plugins + API · 包管理安装 · 编程式调用"),
]

for i, (num, title, desc) in enumerate(layers):
    y = Inches(2.6 + i * 1.5)
    add_line(s, Inches(0.6), y, Inches(0.08), Inches(1.2), IKB)  # left accent bar
    add_rect(s, Inches(0.9), y, Inches(11.8), Inches(1.2), fill_color=GREY_1)
    add_textbox(s, Inches(1.2), y + Inches(0.1), Inches(1.0), Inches(0.5),
                num, Pt(28), "Microsoft YaHei", IKB, line_spacing=1.0)
    add_textbox(s, Inches(2.2), y + Inches(0.1), Inches(4.0), Inches(0.35),
                title, Pt(18), "Microsoft YaHei", INK)
    add_textbox(s, Inches(2.2), y + Inches(0.5), Inches(10.0), Inches(0.6),
                desc, Pt(11), "Microsoft YaHei", GREY_3, line_spacing=1.4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 · S15 Feature Matrix
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, PAPER)
chrome_text(s, "Feature Matrix · 功能矩阵", "08 / 10", 8)

add_textbox(s, Inches(0.6), Inches(1.0), Inches(6), Inches(0.25),
            "FEATURES AT A GLANCE", Pt(10), "Consolas", IKB, bold=True)
add_textbox(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
            "功能覆盖一览", Pt(36), "Microsoft YaHei", INK, line_spacing=0.95)

features = [
    ("PLANNING", "Plan Mode", "Claude ✓"),
    ("EXECUTION", "Agent Loop", "Codex ✓"),
    ("MEMORY", "跨会话记忆", "Claude ✓"),
    ("ISOLATION", "Sandbox", "Codex ✓"),
    ("SKILLS", "Markdown Skills", "Claude ✓"),
    ("MULTIMODAL", "多模态理解", "Codex ✓"),
    ("HOOKS", "事件驱动钩子", "Claude ✓"),
    ("WEB", "Web 搜索取回", "Claude ✓"),
]

positions = [
    (Inches(0.6), Inches(2.2)), (Inches(3.7), Inches(2.2)),
    (Inches(6.8), Inches(2.2)), (Inches(9.9), Inches(2.2)),
    (Inches(0.6), Inches(3.9)), (Inches(3.7), Inches(3.9)),
    (Inches(6.8), Inches(3.9)), (Inches(9.9), Inches(3.9)),
]

for (cat, feat, tag), (x, y) in zip(features, positions):
    add_rect(s, x, y, Inches(2.8), Inches(1.3), fill_color=GREY_1)
    add_textbox(s, x + Inches(0.2), y + Inches(0.1), Inches(2.4), Inches(0.2),
                cat, Pt(8), "Consolas", GREY_3)
    add_textbox(s, x + Inches(0.2), y + Inches(0.35), Inches(2.4), Inches(0.35),
                feat, Pt(14), "Microsoft YaHei", INK)
    add_textbox(s, x + Inches(0.2), y + Inches(0.85), Inches(2.4), Inches(0.25),
                tag, Pt(9), "Consolas", IKB)

# Bottom stat
add_line(s, Inches(0.6), Inches(5.5), Inches(12), Pt(1), GREY_2)
add_textbox(s, Inches(0.6), Inches(5.7), Inches(6), Inches(1.0),
            "各有千秋", Pt(40), "Microsoft YaHei", IKB, line_spacing=0.95)
add_textbox(s, Inches(6.5), Inches(6.1), Inches(6), Inches(0.3),
            "没有全能的工具 · 只有合适的工具", Pt(12), "Consolas", GREY_3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 · S12 Manifesto + Banner (Dark)
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, INK)
chrome_text(s, "Conclusion · 结论", "09 / 10", 9)

add_textbox(s, Inches(0.6), Inches(1.2), Inches(6), Inches(0.25),
            "THE BOTTOM LINE", Pt(10), "Consolas", RGBColor(0xAA, 0xAA, 0xAA))

add_textbox(s, Inches(0.6), Inches(1.6), Inches(10), Inches(1.2),
            "选工具\n不是选模型", Pt(44), "Microsoft YaHei", WHITE,
            line_spacing=0.94)

tf = add_rich_textbox(s, Inches(0.6), Inches(3.2), Inches(10), Inches(1.5))
add_para(tf,
    "模型能力差距正在缩小,真正的差异化在于工具链设计哲学和工作流适配度。"
    "Claude Code 适合需要精细控制和安全保障的团队;"
    "Codex 适合追求自动化程度和云原生体验的团队。",
    Pt(16), "Microsoft YaHei", RGBColor(0xCC, 0xCC, 0xCC), line_spacing=1.5)

# Bottom accent banner
add_rect(s, Inches(0), Inches(5.8), W, Inches(1.7), fill_color=IKB)
add_textbox(s, Inches(0.6), Inches(6.0), Inches(11), Inches(0.8),
            "理解差异,才能做出正确选择", Pt(32), "Microsoft YaHei", WHITE,
            line_spacing=0.95)
add_textbox(s, Inches(0.6), Inches(6.8), Inches(11), Inches(0.3),
            "UNDERSTAND THE DIFFERENCE, MAKE THE RIGHT CHOICE",
            Pt(10), "Consolas", RGBColor(0xCC, 0xCC, 0xFF))

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 · S10 Split Closing
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
# Left half - IKB
add_rect(s, Inches(0), Inches(0), Inches(6.5), H, fill_color=IKB)
# Right half - White
add_rect(s, Inches(6.5), Inches(0), Inches(6.833), H, fill_color=PAPER)

# Left side content
add_textbox(s, Inches(0.5), Inches(0.35), Inches(5.5), Inches(0.3),
            "10 / 10", Pt(9), "Consolas", RGBColor(0xCC, 0xCC, 0xFF),
            alignment=PP_ALIGN.LEFT)
add_textbox(s, Inches(0.5), Inches(0.35), Inches(5.5), Inches(0.3),
            "CLOSING", Pt(9), "Consolas", RGBColor(0xCC, 0xCC, 0xFF),
            alignment=PP_ALIGN.RIGHT)

add_textbox(s, Inches(0.5), Inches(1.8), Inches(5.5), Inches(0.25),
            "TAKEAWAY", Pt(10), "Consolas", RGBColor(0xCC, 0xCC, 0xFF))
add_textbox(s, Inches(0.5), Inches(2.2), Inches(5.5), Inches(2.5),
            "Choose\nthe workflow\nnot the tool", Pt(40), "Microsoft YaHei",
            WHITE, line_spacing=0.94)
add_textbox(s, Inches(0.5), Inches(5.0), Inches(5.5), Inches(0.8),
            "工具为你服务,而非你为工具服务.\n理解差异,根据团队工作流做出选择.",
            Pt(13), "Microsoft YaHei", RGBColor(0xCC, 0xCC, 0xFF),
            line_spacing=1.5)

add_line(s, Inches(0.5), Inches(6.2), Inches(5.5), Pt(1),
         RGBColor(0x55, 0x55, 0xCC))
add_textbox(s, Inches(0.5), Inches(6.4), Inches(5.5), Inches(0.3),
            "内部技术分享  ·  2026.05.13", Pt(10), "Consolas",
            RGBColor(0xAA, 0xAA, 0xDD))

# Right side content
add_textbox(s, Inches(7.0), Inches(0.35), Inches(5.5), Inches(0.3),
            "TAKEAWAYS", Pt(9), "Consolas", GREY_3)
add_textbox(s, Inches(7.0), Inches(0.35), Inches(5.5), Inches(0.3),
            "03 RULES", Pt(9), "Consolas", GREY_3, alignment=PP_ALIGN.RIGHT)

takeaways = [
    ("01", "先理解工作流,再选工具",
     "你的团队需要精细控制还是高度自动化?\n答案决定了应该选谁."),
    ("02", "关注生态,不只是模型",
     "Skills / MCP / Hooks 的生态网络,\n比单次 benchmark 跑分更重要."),
    ("03", "两者可以共存",  # accent
     "不同场景用不同工具:\nClaude Code 做架构设计, Codex 做快速原型."),
]

for i, (num, title, desc) in enumerate(takeaways):
    y = Inches(1.5 + i * 1.8)
    is_accent = (num == "03")
    num_color = IKB if is_accent else INK
    title_color = IKB if is_accent else INK

    add_line(s, Inches(7.0), y, Inches(5.5), Pt(1), GREY_2)
    add_textbox(s, Inches(7.0), y + Inches(0.15), Inches(0.6), Inches(0.5),
                num, Pt(24), "Microsoft YaHei", num_color, line_spacing=0.9)
    add_textbox(s, Inches(7.8), y + Inches(0.15), Inches(4.5), Inches(0.35),
                title, Pt(18), "Microsoft YaHei", title_color)
    add_textbox(s, Inches(7.8), y + Inches(0.55), Inches(4.5), Inches(0.7),
                desc, Pt(11), "Microsoft YaHei", GREY_3, line_spacing=1.4)

# Bottom accent line on last item
if is_accent:
    add_line(s, Inches(7.0), Inches(1.5 + 2 * 1.8 + 1.5), Inches(5.5), Pt(2), IKB)

add_textbox(s, Inches(7.0), Inches(6.8), Inches(5.5), Inches(0.3),
            "→ 完 · END OF DECK", Pt(10), "Consolas", GREY_3,
            alignment=PP_ALIGN.RIGHT)

# ── Save ──
output_path = r"D:\claude code\ppt\Claude_Code_vs_Codex.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"   {len(prs.slides)} slides - 16:9 - IKB Swiss Style")
